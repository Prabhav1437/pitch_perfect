
import json
import logging
import re
from typing import Dict, List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


from gemini_evaluator import GeminiEvaluator
from llm_evaluator import LLMEvaluator

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTReconstructor:
    def __init__(self, gemini_evaluator=None, llm_evaluator=None, output_dir: str = "generated_ppts"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        # Use Gemini primarily if available
        self.evaluator = gemini_evaluator or GeminiEvaluator()
        # Fallback to local model if strict dependency exists, but we prefer Gemini for structure gen
        self.fallback_evaluator = llm_evaluator or LLMEvaluator()

    def reconstruct_presentation(
        self, 
        summary: str, 
        problem_statement: str, 
        analysis: Dict, 
        custom_instructions: str = ""
    ) -> Dict:
        """
        Reconstruct a presentation based on analysis and return the path and structure.
        """
        structure = self.generate_structure(summary, problem_statement, analysis, custom_instructions)
        
        if not structure or "slides" not in structure:
            logger.warning("Failed to generate valid structure. Using fallback.")
            structure = self._get_fallback_structure()
            
        filename = f"reconstructed_{hash(summary[:20])}.pptx"
        file_path = self.create_ppt(structure, filename)
        
        return {
            "file_path": file_path,
            "structure": structure
        }

    def generate_structure(
        self, 
        summary: str, 
        problem_statement: str, 
        analysis: Dict,
        custom_instructions: str
    ) -> Dict:
        logger.info("Generating PPT structure using Gemini/AI...")
        
        prompt = f"""
        Act as a professional pitch deck designer and content strategist. 
        Your task is to RECONSTRUCT and IMPROVE a pitch deck based on the following project summary and judge's critique.
        
        CONTEXT:
        Project Summary: {summary}
        Problem Statement to Address: {problem_statement}
        
        CRITIQUE (Fix these issues):
        Weaknesses: {json.dumps(analysis.get('weaknesses', []), indent=2)}
        Missing Elements: {json.dumps(analysis.get('missing_elements', []), indent=2)}
        Detailed Analysis: {json.dumps(analysis.get('detailed_analysis', {}), indent=2)}
        
        USER INSTRUCTIONS:
        {custom_instructions if custom_instructions else "Improve the flow, clarity, and impact."}
        
        TASK:
        Generate a JSON structure for a complete, professional pitch deck (10-12 slides).
        Each slide must have a 'title', 'layout', and 'content'.
        
        LAYOUT OPTIONS:
        - "Title Slide" (Title + Subtitle)
        - "Title and Content" (Title + Bullet points)
        - "Two Content" (Title + Left Bullets + Right Bullets)
        - "Section Header" (Title only)
        
        OUTPUT FORMAT (JSON ONLY, NO MARKDOWN):
        {{
            "slides": [
                {{
                    "title": "Slide Title",
                    "layout": "Title Slide",
                    "content": {{
                        "title": "Main Title",
                        "subtitle": "Subtitle/Tagline"
                    }}
                }},
                {{
                    "title": "The Problem",
                    "layout": "Title and Content",
                    "content": {{
                        "bullets": [
                            "Point 1...",
                            "Point 2..."
                        ]
                    }}
                }}
            ]
        }}
        """
        
        # Try primary evaluator (Gemini)
        response_text = self.evaluator.generate_response(prompt)
        
        # If Gemini fails (e.g. key missing) and returns error string, maybe fallback?
        # But user specifically asked for Gemini.
        # We assume extract_json handles the string parsing
        structure = self.evaluator.extract_json_from_text(response_text)
        
        if not structure:
             logger.warning("Gemini failed to return valid JSON. Attempting fallback to local LLM.")
             try:
                 # Attempt fallback if Gemini failed completely
                 response_text = self.fallback_evaluator.generate_response(prompt, max_new_tokens=2048)
                 structure = self.fallback_evaluator.extract_json_from_text(response_text)
             except Exception as e:
                 logger.error(f"Fallback failed: {e}")
        
        return structure

    def create_ppt(self, structure: Dict, filename: str = "reconstructed.pptx") -> str:
        logger.info(f"Creating PPT file: {filename}")
        prs = Presentation()
        
        # Helper to set text frame content
        def set_text_frame(shape, text_data):
            if not shape.has_text_frame:
                return
            text_frame = shape.text_frame
            text_frame.clear()  # Clear existing
            
            if isinstance(text_data, str):
                p = text_frame.paragraphs[0]
                p.text = text_data
            elif isinstance(text_data, list):
                for i, line in enumerate(text_data):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = str(line)
                    p.space_before = Pt(6)

        for slide_data in structure.get("slides", []):
            layout_name = slide_data.get("layout", "Title and Content")
            title_text = slide_data.get("title", "Untitled")
            content_data = slide_data.get("content", {})
            
            # Map layout to index
            # 0: Title, 1: Title & Content, 2: Section Header, 3: Two Content
            layout_index = 1
            if "Title Slide" in layout_name: layout_index = 0
            elif "Section Header" in layout_name: layout_index = 2
            elif "Two Content" in layout_name: layout_index = 3
            
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # Fill placeholders based on layout
            if layout_index == 0: # Title Slide
                if len(slide.placeholders) > 1:
                    set_text_frame(slide.placeholders[1], content_data.get("subtitle", ""))
                # Main title is already set above
                if slide.shapes.title and "title" in content_data:
                     slide.shapes.title.text = content_data["title"]

            elif layout_index == 1: # Title and Content
                if len(slide.placeholders) > 1:
                    # Try 'bullets' or 'body' or 'text'
                    body_text = content_data.get("bullets") or content_data.get("body") or content_data.get("text", "")
                    set_text_frame(slide.placeholders[1], body_text)
                    
            elif layout_index == 3: # Two Content
                # placeholders[1] is left, placeholders[2] is right usually
                if len(slide.placeholders) > 1:
                    set_text_frame(slide.placeholders[1], content_data.get("left", []))
                if len(slide.placeholders) > 2:
                    set_text_frame(slide.placeholders[2], content_data.get("right", []))

        output_path = self.output_dir / filename
        prs.save(str(output_path))
        logger.info(f"PPT saved to {output_path}")
        return str(output_path)

    def _get_fallback_structure(self) -> Dict:
        return {
            "slides": [
                {
                    "title": "Project Title",
                    "layout": "Title Slide",
                    "content": {"title": "Project Name", "subtitle": "Innovative Solution"}
                },
                {
                    "title": "Problem",
                    "layout": "Title and Content",
                    "content": {"bullets": ["State the problem specifically", "Explain the pain point"]}
                },
                {
                    "title": "Solution",
                    "layout": "Title and Content",
                    "content": {"bullets": ["Describe your solution", "Highlight key benefits"]}
                }
            ]
        }

    def refine_structure(
        self,
        current_structure: Dict,
        user_instruction: str,
        presentation_summary: str
    ) -> Dict:
        """
        Update the PPT structure based on user chat instruction.
        """
        logger.info("Refining PPT structure based on user feedback...")
        
        prompt = f"""
        Act as a professional pitch deck designer.
        The user wants to modify an existing presentation structure.
        
        CONTEXT:
        Project Summary: {presentation_summary}
        Current Structure (JSON):
        {json.dumps(current_structure, indent=2)}
        
        USER INSTRUCTION:
        "{user_instruction}"
        
        TASK:
        Modify the JSON structure to address the user's request.
        - If they ask to add a slide, insert it in the appropriate position.
        - If they ask to change text, update the content.
        - Maintain the JSON format exactly.
        
        OUTPUT FORMAT (JSON ONLY, NO MARKDOWN):
        {{
            "slides": [...]
        }}
        """
        
        response_text = self.evaluator.generate_response(prompt)
        new_structure = self.evaluator.extract_json_from_text(response_text)
        
        if not new_structure or "slides" not in new_structure:
            logger.warning("Failed to parse refined structure. Returning original.")
            # Be strict about failure here to let user know
            return {
                "structure": current_structure,
                "file_path": None 
            }
            
        # Re-generate PPT file
        filename = f"refined_{hash(user_instruction[:20])}.pptx"
        file_path = self.create_ppt(new_structure, filename)
        
        return {
            "structure": new_structure,
            "file_path": file_path
        }

if __name__ == "__main__":
    # Test run
    reconstructor = PPTReconstructor()
    structure = reconstructor.generate_structure(
        "A waste management AI system.",
        "Solve recycling inefficiency.",
        {"weaknesses": ["No business model"], "missing_elements": ["Revenue stream"]},
        ""
    )
    print(json.dumps(structure, indent=2))
    reconstructor.create_ppt(structure, "test_reconstruct.pptx")
