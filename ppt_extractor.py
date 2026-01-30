"""
PowerPoint extraction module.
Extracts text content from PPTX files including titles, body text, and speaker notes.
"""
from pptx import Presentation
from typing import List, Dict, Optional
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTExtractor:
    """Extract structured content from PowerPoint presentations."""
    
    def __init__(self):
        """Initialize the PPT extractor."""
        pass
    
    def extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a shape object.
        
        Args:
            shape: A shape object from python-pptx
            
        Returns:
            Extracted text as string
        """
        text_content = []
        
        if hasattr(shape, "text") and shape.text:
            text_content.append(shape.text.strip())
        
        # Handle tables
        if hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_content.append(cell.text.strip())
        
        return " ".join(text_content)
    
    def extract_slide_content(self, slide) -> Dict[str, any]:
        """
        Extract all content from a single slide.
        
        Args:
            slide: A slide object from python-pptx
            
        Returns:
            Dictionary containing slide content
        """
        slide_data = {
            "title": "",
            "content": [],
            "notes": ""
        }
        
        # Extract title
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()
        
        # Extract body content
        for shape in slide.shapes:
            # Skip title shape as we already extracted it
            if shape == slide.shapes.title:
                continue
            
            text = self.extract_text_from_shape(shape)
            if text:
                slide_data["content"].append(text)
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_data["notes"] = notes_slide.notes_text_frame.text.strip()
        
        return slide_data
    
    def extract_from_file(self, file_path: str) -> Dict[str, any]:
        """
        Extract all content from a PowerPoint file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Dictionary containing presentation metadata and slide content
            
        Raises:
            Exception: If file cannot be read or is corrupted
        """
        try:
            logger.info(f"Extracting content from: {file_path}")
            prs = Presentation(file_path)
            
            presentation_data = {
                "slide_count": len(prs.slides),
                "slides": []
            }
            
            for idx, slide in enumerate(prs.slides, 1):
                slide_content = self.extract_slide_content(slide)
                slide_content["slide_number"] = idx
                presentation_data["slides"].append(slide_content)
                
                logger.debug(f"Extracted slide {idx}: {slide_content['title'][:50]}")
            
            logger.info(f"Successfully extracted {len(prs.slides)} slides")
            return presentation_data
            
        except Exception as e:
            logger.error(f"Error extracting PPT content: {str(e)}")
            raise Exception(f"Failed to extract PowerPoint content: {str(e)}")
    
    def get_full_text(self, presentation_data: Dict[str, any]) -> str:
        """
        Combine all text from presentation into a single string.
        
        Args:
            presentation_data: Dictionary from extract_from_file()
            
        Returns:
            Combined text from all slides
        """
        full_text_parts = []
        
        for slide in presentation_data["slides"]:
            slide_parts = []
            
            if slide["title"]:
                slide_parts.append(f"Slide {slide['slide_number']}: {slide['title']}")
            
            if slide["content"]:
                slide_parts.append(" ".join(slide["content"]))
            
            if slide["notes"]:
                slide_parts.append(f"Notes: {slide['notes']}")
            
            if slide_parts:
                full_text_parts.append("\n".join(slide_parts))
        
        return "\n\n".join(full_text_parts)


if __name__ == "__main__":
    # Example usage
    extractor = PPTExtractor()
    # data = extractor.extract_from_file("sample.pptx")
    # print(f"Extracted {data['slide_count']} slides")
