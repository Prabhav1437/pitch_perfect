from pptx import Presentation

prs = Presentation()

# Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "EcoSort AI"
subtitle.text = "Revolutionizing Waste Management with Computer Vision"

# Problem Slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Problem"
content = slide.placeholders[1]
content.text = "Global waste is growing by 20% annually.\nRecycling centers are slow and inefficient.\nHuman sorting is dangerous and error-prone."

# Solution Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Solution"
content = slide.placeholders[1]
content.text = "EcoSort AI uses cameras and robotic arms to sort waste.\n99% accuracy using Deep Learning.\nSorts 10x faster than humans."

# Tech Stack Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Tech Stack"
content = slide.placeholders[1]
content.text = "Frontend: React Native App\nBackend: Python & FastAPI\nAI: PyTorch & YOLOv8\nHardware: Raspberry Pi & Servo Motors"

# Business Model Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Business Model"
content = slide.placeholders[1]
content.text = "B2B SaaS Subscription for Analytics.\nOne-time hardware cost.\n$50k ARR projected in Year 1."

prs.save('/Users/adityasmac/Desktop/ppt evaluator/test_deck.pptx')
print("Created test_deck.pptx")
