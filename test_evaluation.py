"""
Integration tests for the evaluation system.
"""
import pytest
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

from ppt_extractor import PPTExtractor
from summarizer import Summarizer
from semantic_scorer import SemanticScorer
from llm_evaluator import LLMEvaluator
from evaluation_orchestrator import EvaluationOrchestrator
from models import EvaluationResponse


class TestPPTExtractor:
    """Test PPT extraction functionality."""
    
    def test_extractor_initialization(self):
        """Test that extractor initializes correctly."""
        extractor = PPTExtractor()
        assert extractor is not None
    
    def test_create_sample_ppt(self, tmp_path):
        """Create a sample PPT for testing."""
        prs = Presentation()
        
        # Slide 1: Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Web Scraping Project"
        slide1.placeholders[1].text = "A comprehensive solution"
        
        # Slide 2: Content slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Technologies Used"
        content = slide2.placeholders[1].text_frame
        content.text = "Python"
        content.add_paragraph().text = "BeautifulSoup"
        content.add_paragraph().text = "Requests library"
        
        # Slide 3: Implementation
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "Implementation Details"
        content = slide3.placeholders[1].text_frame
        content.text = "We built a modular scraper that handles dynamic content"
        
        # Save
        ppt_path = tmp_path / "test_presentation.pptx"
        prs.save(str(ppt_path))
        
        return str(ppt_path)
    
    def test_extract_from_file(self, tmp_path):
        """Test extracting content from a PPT file."""
        ppt_path = self.test_create_sample_ppt(tmp_path)
        
        extractor = PPTExtractor()
        data = extractor.extract_from_file(ppt_path)
        
        assert data is not None
        assert 'slide_count' in data
        assert data['slide_count'] == 3
        assert 'slides' in data
        assert len(data['slides']) == 3
        
        # Check first slide
        assert data['slides'][0]['title'] == "Web Scraping Project"
    
    def test_get_full_text(self, tmp_path):
        """Test getting full text from presentation."""
        ppt_path = self.test_create_sample_ppt(tmp_path)
        
        extractor = PPTExtractor()
        data = extractor.extract_from_file(ppt_path)
        full_text = extractor.get_full_text(data)
        
        assert "Web Scraping" in full_text
        assert "Python" in full_text
        assert "BeautifulSoup" in full_text


class TestSummarizer:
    """Test summarization functionality."""
    
    def test_summarizer_initialization(self):
        """Test that summarizer initializes correctly."""
        summarizer = Summarizer()
        assert summarizer is not None
        assert summarizer.model_name is not None
    
    def test_summarize_short_text(self):
        """Test that short text is not summarized."""
        summarizer = Summarizer()
        short_text = "This is a short text."
        result = summarizer.summarize_text(short_text)
        assert result == short_text


class TestSemanticScorer:
    """Test semantic scoring functionality."""
    
    def test_scorer_initialization(self):
        """Test that scorer initializes correctly."""
        scorer = SemanticScorer()
        assert scorer is not None
        assert scorer.model_name is not None
    
    def test_cosine_similarity(self):
        """Test cosine similarity calculation."""
        import numpy as np
        scorer = SemanticScorer()
        
        # Identical vectors should have similarity 1.0
        vec1 = np.array([1, 0, 0])
        vec2 = np.array([1, 0, 0])
        similarity = scorer.cosine_similarity(vec1, vec2)
        assert abs(similarity - 1.0) < 0.01
        
        # Orthogonal vectors should have similarity 0.0
        vec3 = np.array([1, 0, 0])
        vec4 = np.array([0, 1, 0])
        similarity = scorer.cosine_similarity(vec3, vec4)
        assert abs(similarity - 0.0) < 0.01


class TestLLMEvaluator:
    """Test LLM evaluation functionality."""
    
    def test_evaluator_initialization(self):
        """Test that evaluator initializes correctly."""
        evaluator = LLMEvaluator()
        assert evaluator is not None
        assert evaluator.model_name is not None
    
    def test_extract_json_from_text(self):
        """Test JSON extraction from text."""
        evaluator = LLMEvaluator()
        
        text = 'Some text before {"key": "value", "number": 42} some text after'
        result = evaluator.extract_json_from_text(text)
        
        assert result is not None
        assert result['key'] == 'value'
        assert result['number'] == 42
    
    def test_validate_and_fix_response(self):
        """Test response validation and fixing."""
        evaluator = LLMEvaluator()
        
        # Incomplete response
        incomplete = {
            "scores": {
                "relevance": 8
            }
        }
        
        fixed = evaluator.validate_and_fix_response(incomplete)
        
        assert 'scores' in fixed
        assert 'relevance' in fixed['scores']
        assert 'clarity' in fixed['scores']
        assert 'overall_score' in fixed
        assert 'strengths' in fixed
        assert isinstance(fixed['strengths'], list)
    
    def test_get_default_evaluation(self):
        """Test default evaluation generation."""
        evaluator = LLMEvaluator()
        default = evaluator.get_default_evaluation()
        
        assert 'scores' in default
        assert 'overall_score' in default
        assert default['overall_score'] == 25
        assert 'strengths' in default
        assert 'weaknesses' in default


class TestEvaluationResponse:
    """Test Pydantic models."""
    
    def test_valid_response(self):
        """Test creating a valid response."""
        data = {
            "scores": {
                "relevance": 8.5,
                "clarity": 7.0,
                "technical_accuracy": 9.0,
                "structure": 8.0,
                "completeness": 7.5
            },
            "overall_score": 40.0,
            "strengths": ["Good technical depth", "Clear structure", "Comprehensive coverage"],
            "weaknesses": ["Could improve visuals", "Missing some examples"],
            "improvement_suggestions": ["Add more diagrams", "Include code examples"],
            "missing_elements": ["Performance metrics"],
            "summary_evaluation": "Strong presentation with good technical content."
        }
        
        response = EvaluationResponse(**data)
        assert response.overall_score == 40.0
        assert len(response.strengths) == 3
    
    def test_score_validation(self):
        """Test that scores are validated."""
        data = {
            "scores": {
                "relevance": 8.0,
                "clarity": 7.0,
                "technical_accuracy": 9.0,
                "structure": 8.0,
                "completeness": 8.0
            },
            "overall_score": 35.0,  # Should be 40, will be auto-corrected
            "strengths": ["Good"],
            "weaknesses": ["Bad"],
            "improvement_suggestions": ["Improve"],
            "missing_elements": [],
            "summary_evaluation": "Test evaluation"
        }
        
        response = EvaluationResponse(**data)
        # Should auto-correct to sum of scores
        assert response.overall_score == 40.0


# Integration test
def test_full_pipeline(tmp_path):
    """Test the complete evaluation pipeline."""
    # Create sample PPT
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Web Scraping with Python"
    slide.placeholders[1].text = "Using BeautifulSoup and Requests"
    
    ppt_path = tmp_path / "test.pptx"
    prs.save(str(ppt_path))
    
    # Note: This test requires models to be downloaded
    # Uncomment to run full integration test
    """
    orchestrator = EvaluationOrchestrator()
    result = orchestrator.evaluate_presentation(
        str(ppt_path),
        "Build a web scraper for e-commerce websites"
    )
    
    assert 'scores' in result
    assert 'overall_score' in result
    assert result['overall_score'] >= 0
    assert result['overall_score'] <= 50
    """


if __name__ == "__main__":
    # Run tests
    pytest.main([__file__, "-v"])
